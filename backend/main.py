from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Header
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import os
import uuid
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import faiss
import numpy as np
from google.generativeai import GenerativeModel
import google.generativeai as genai
from sentence_transformers import SentenceTransformer
from io import StringIO, BytesIO
from typing import List, Optional, Dict
import logging
import asyncpg
import bcrypt
import jwt
from datetime import datetime, timedelta
import json

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Agentic RAG Chatbot API", version="1.0.0")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://127.0.0.1:5173"],
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type", "Accept", "X-Requested-With"],
)

# Database configuration
DATABASE_URL = "postgresql://admin:password@localhost:5432/rag_chatbot1"

# JWT configuration
SECRET_KEY = "jwt_secert_key"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

# Configure Gemini API
genai.configure(api_key="your_gemini_api_key")

# Initialize Gemini model
model = GenerativeModel('gemini-2.5-flash')

# Initialize SentenceTransformer for embeddings
embedder = SentenceTransformer('all-MiniLM-L6-v2')

# Initialize FAISS index
dimension = 384
index = faiss.IndexFlatL2(dimension)

# Global stores
chunks = []
chunk_embeddings = []

# Pydantic models
class QueryRequest(BaseModel):
    query: str
    session_id: str

class QueryResponse(BaseModel):
    answer: str
    sources: List[str]
    query: str
    session_id: str

class ChatHistoryItem(BaseModel):
    id: str
    query: str
    answer: str
    sources: List[str]
    timestamp: datetime

class ChatSession(BaseModel):
    id: str
    user_id: str
    filename: str
    created_at: datetime
    messages: List[ChatHistoryItem]

class UploadResponse(BaseModel):
    message: str
    filename: str
    chunks_count: int
    session_id: str

class UserCreate(BaseModel):
    name: str
    email: str
    password: str

class UserLogin(BaseModel):
    email: str
    password: str

class UserResponse(BaseModel):
    id: str
    name: str
    email: str

# Database setup
async def get_db():
    try:
        conn = await asyncpg.connect(DATABASE_URL)
        yield conn
    except Exception as e:
        logger.error(f"Database connection error: {str(e)}")
        raise HTTPException(status_code=500, detail="Database connection failed")
    finally:
        if 'conn' in locals():
            await conn.close()

# JWT authentication
async def get_current_user(authorization: Optional[str] = Header(None), db: asyncpg.Connection = Depends(get_db)):
    if not authorization:
        logger.error("No Authorization header provided")
        raise HTTPException(status_code=401, detail="Authorization header missing. Please log in.")
    
    try:
        token = authorization.replace("Bearer ", "")
        logger.info(f"Decoding token: {token[:10]}...")
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email = payload.get("sub")
        if not email:
            logger.error("No 'sub' claim in token")
            raise HTTPException(status_code=401, detail="Invalid token. Please log in.")
        
        user = await db.fetchrow("SELECT id, name, email FROM users WHERE email = $1", email)
        if not user:
            logger.error(f"User not found for email: {email}")
            raise HTTPException(status_code=401, detail="User not found. Please log in.")
        
        logger.info(f"Authenticated user: {email}")
        return UserResponse(id=str(user["id"]), name=user["name"], email=user["email"])
    except jwt.ExpiredSignatureError:
        logger.error("Token expired")
        raise HTTPException(status_code=401, detail="Token expired. Please log in again.")
    except jwt.InvalidTokenError as e:
        logger.error(f"Invalid token: {str(e)}")
        raise HTTPException(status_code=401, detail="Invalid token. Please log in.")
    except Exception as e:
        logger.error(f"JWT decode error: {str(e)}")
        raise HTTPException(status_code=401, detail="Invalid token. Please log in.")

class MCPMessage:
    def __init__(self, sender, receiver, type, payload, trace_id=None):
        self.sender = sender
        self.receiver = receiver
        self.type = type
        self.trace_id = trace_id or str(uuid.uuid4())
        self.payload = payload

    def to_dict(self):
        return {
            "sender": self.sender,
            "receiver": self.receiver,
            "type": self.type,
            "trace_id": self.trace_id,
            "payload": self.payload
        }

class IngestionAgent:
    def process_document(self, file_content: bytes, filename: str):
        extension = os.path.splitext(filename)[1].lower()
        text = ""
        
        try:
            if extension == '.pdf':
                pdf_file = BytesIO(file_content)
                reader = PdfReader(pdf_file)
                for page in reader.pages:
                    text += page.extract_text() or ""
            elif extension == '.docx':
                doc_file = BytesIO(file_content)
                doc = Document(doc_file)
                text = "\n".join([para.text for para in doc.paragraphs])
            elif extension == '.pptx':
                pptx_file = BytesIO(file_content)
                prs = Presentation(pptx_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif extension == '.csv':
                csv_file = StringIO(file_content.decode('utf-8'))
                df = pd.read_csv(csv_file)
                text = df.to_string()
            elif extension in ['.txt', '.md']:
                text = file_content.decode("utf-8")
        except Exception as e:
            logger.error(f"Error processing {filename}: {str(e)}")
            raise HTTPException(status_code=400, detail=f"Error processing {filename}: {str(e)}")

        chunk_size = 500
        chunk_list = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
        
        message = MCPMessage(
            sender="IngestionAgent",
            receiver="RetrievalAgent",
            type="DOCUMENT_PROCESSED",
            payload={
                "filename": filename,
                "chunks": chunk_list
            }
        )
        return message

class RetrievalAgent:
    def store_chunks(self, message, session_id: str):
        global chunks, chunk_embeddings, index
        chunk_list = message.payload["chunks"]
        filename = message.payload["filename"]
        
        embeddings = embedder.encode(chunk_list)
        
        for chunk, embedding in zip(chunk_list, embeddings):
            chunks.append({"text": chunk, "filename": filename, "session_id": session_id})
            chunk_embeddings.append(embedding)
        
        index.add(np.array(embeddings, dtype=np.float32))
        
        return MCPMessage(
            sender="RetrievalAgent",
            receiver="CoordinatorAgent",
            type="CHUNKS_STORED",
            payload={"status": "success", "filename": filename, "chunks_count": len(chunk_list), "session_id": session_id}
        )

    def retrieve(self, message):
        if len(chunks) == 0:
            return MCPMessage(
                sender="RetrievalAgent",
                receiver="LLMResponseAgent",
                type="CONTEXT_RESPONSE",
                payload={
                    "query": message.payload["query"],
                    "session_id": message.payload["session_id"],
                    "top_chunks": []
                }
            )
            
        query = message.payload["query"]
        session_id = message.payload["session_id"]
        query_embedding = embedder.encode([query])[0]
        
        k = min(3, len(chunks))
        distances, indices = index.search(np.array([query_embedding], dtype=np.float32), k)
        
        relevant_chunks = [chunks[i] for i in indices[0] if i < len(chunks) and chunks[i]["session_id"] == session_id]
        
        return MCPMessage(
            sender="RetrievalAgent",
            receiver="LLMResponseAgent",
            type="CONTEXT_RESPONSE",
            payload={
                "query": query,
                "session_id": session_id,
                "top_chunks": relevant_chunks
            }
        )

class LLMResponseAgent:
    def generate_response(self, message):
        query = message.payload["query"]
        session_id = message.payload["session_id"]
        context = "\n".join([chunk["text"] for chunk in message.payload["top_chunks"]])
        
        if not context.strip():
            prompt = f"""
            Question: {query}
            
            I don't have any relevant documents uploaded to answer this question. Please upload some documents first and then ask your question.
            """
        else:
            prompt = f"""
            Context: {context}
            
            Question: {query}
            
            Provide a concise and accurate answer based on the context. Include source information if relevant.
            """
        
        try:
            response = model.generate_content(prompt)
            answer = response.text
            
            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="FINAL_RESPONSE",
                payload={
                    "query": query,
                    "answer": answer,
                    "sources": list(set([chunk["filename"] for chunk in message.payload["top_chunks"]])),
                    "session_id": session_id
                }
            )
        except Exception as e:
            logger.error(f"Error generating response: {str(e)}")
            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                payload={"error": str(e), "session_id": session_id}
            )

class CoordinatorAgent:
    def __init__(self):
        self.ingestion_agent = IngestionAgent()
        self.retrieval_agent = RetrievalAgent()
        self.llm_response_agent = LLMResponseAgent()

    async def process_upload(self, file_content: bytes, filename: str, user_id: str, db: asyncpg.Connection):
        message = self.ingestion_agent.process_document(file_content, filename)
        if message:
            session_id = str(uuid.uuid4())
            await db.execute(
                """
                INSERT INTO chat_sessions (id, user_id, filename, created_at)
                VALUES ($1, $2, $3, $4)
                """,
                session_id, user_id, filename, datetime.utcnow()
            )
            return self.retrieval_agent.store_chunks(message, session_id)
        return None

    def process_query(self, query: str, session_id: str):
        message = MCPMessage(
            sender="CoordinatorAgent",
            receiver="RetrievalAgent",
            type="QUERY",
            payload={"query": query, "session_id": session_id}
        )
        retrieval_message = self.retrieval_agent.retrieve(message)
        return self.llm_response_agent.generate_response(retrieval_message)

    def start_new_session(self):
        global chunks, chunk_embeddings, index
        chunks = []
        chunk_embeddings = []
        index = faiss.IndexFlatL2(dimension)
        return {"message": "New session started, context cleared"}

# Initialize coordinator
coordinator = CoordinatorAgent()

@app.get("/")
async def root():
    return {"message": "Agentic RAG Chatbot API is running!"}

@app.post("/api/auth/signup")
async def signup(user: UserCreate, db: asyncpg.Connection = Depends(get_db)):
    try:
        hashed_password = bcrypt.hashpw(user.password.encode('utf-8'), bcrypt.gensalt())
        user_id = str(uuid.uuid4())
        await db.execute(
            """
            INSERT INTO users (id, name, email, password)
            VALUES ($1, $2, $3, $4)
            """,
            user_id, user.name, user.email, hashed_password
        )
        return {"message": "User created successfully"}
    except asyncpg.UniqueViolationError:
        raise HTTPException(status_code=400, detail="Email already exists")
    except Exception as e:
        logger.error(f"Signup error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/auth/login")
async def login(user: UserLogin, db: asyncpg.Connection = Depends(get_db)):
    try:
        db_user = await db.fetchrow("SELECT id, name, email, password FROM users WHERE email = $1", user.email)
        if not db_user or not bcrypt.checkpw(user.password.encode('utf-8'), db_user["password"]):
            logger.error(f"Invalid login attempt for email: {user.email}")
            raise HTTPException(status_code=401, detail="Invalid email or password")
        
        token = jwt.encode(
            {"sub": db_user["email"], "exp": datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)},
            SECRET_KEY,
            algorithm=ALGORITHM
        )
        
        logger.info(f"Login successful for user: {user.email}")
        return {
            "token": token,
            "user": UserResponse(id=str(db_user["id"]), name=db_user["name"], email=db_user["email"])
        }
    except Exception as e:
        logger.error(f"Login error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/auth/verify")
async def verify_token(current_user: UserResponse = Depends(get_current_user)):
    return {"user": current_user}

@app.post("/upload", response_model=UploadResponse)
async def upload_document(file: UploadFile = File(...), user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        file_content = await file.read()
        message = await coordinator.process_upload(file_content, file.filename, user.id, db)
        
        if message and message.payload["status"] == "success":
            logger.info(f"File uploaded successfully: {file.filename} by user: {user.email}, session_id: {message.payload['session_id']}")
            return UploadResponse(
                message=f"Successfully processed {file.filename}",
                filename=file.filename,
                chunks_count=message.payload["chunks_count"],
                session_id=message.payload["session_id"]
            )
        else:
            raise HTTPException(status_code=400, detail="Failed to process document")
    except Exception as e:
        logger.error(f"Upload error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/api/chat-sessions/new")
async def new_chat_session(user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        session_id = str(uuid.uuid4())
        await db.execute(
            """
            INSERT INTO chat_sessions (id, user_id, filename, created_at)
            VALUES ($1, $2, $3, $4)
            """,
            session_id, user.id, "No file uploaded", datetime.utcnow()
        )
        logger.info(f"New chat session created for user: {user.email}, session_id: {session_id}")
        return {"session_id": session_id, "message": "New chat session created"}
    except Exception as e:
        logger.error(f"New chat session error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))

@app.get("/api/chat-sessions", response_model=List[ChatSession])
async def get_chat_sessions(user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        sessions = await db.fetch(
            """
            SELECT cs.id, cs.user_id, cs.filename, cs.created_at, 
                   COALESCE((
                       SELECT array_agg(
                           json_build_object(
                               'id', m.id,
                               'query', m.query,
                               'answer', m.answer,
                               'sources', m.sources,
                               'timestamp', m.timestamp
                           )
                       )
                       FROM messages m WHERE m.session_id = cs.id
                   ), ARRAY[]::json[]) as messages
            FROM chat_sessions cs
            WHERE cs.user_id = $1
            ORDER BY cs.created_at DESC
            """,
            user.id
        )
        result = []
        for session in sessions:
            # Parse JSON strings in messages array
            messages = [
                ChatHistoryItem(**json.loads(m)) if isinstance(m, str) else ChatHistoryItem(**m)
                for m in session["messages"]
            ]
            result.append(
                ChatSession(
                    id=str(session["id"]),
                    user_id=str(session["user_id"]),
                    filename=session["filename"],
                    created_at=session["created_at"],
                    messages=messages
                )
            )
        logger.info(f"Fetched {len(result)} chat sessions for user: {user.email}")
        return result
    except json.JSONDecodeError as e:
        logger.error(f"JSON parsing error in get_chat_sessions: {str(e)}")
        raise HTTPException(status_code=500, detail="Error parsing chat session data")
    except Exception as e:
        logger.error(f"Get chat sessions error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/chat-sessions/{session_id}", response_model=ChatSession)
async def get_chat_session(session_id: str, user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        session = await db.fetchrow(
            """
            SELECT cs.id, cs.user_id, cs.filename, cs.created_at, 
                   COALESCE((
                       SELECT array_agg(
                           json_build_object(
                               'id', m.id,
                               'query', m.query,
                               'answer', m.answer,
                               'sources', m.sources,
                               'timestamp', m.timestamp
                           )
                       )
                       FROM messages m WHERE m.session_id = cs.id
                   ), ARRAY[]::json[]) as messages
            FROM chat_sessions cs
            WHERE cs.id = $1 AND cs.user_id = $2
            """,
            session_id, user.id
        )
        if not session:
            raise HTTPException(status_code=404, detail="Chat session not found")
        
        # Parse JSON strings in messages array
        messages = [
            ChatHistoryItem(**json.loads(m)) if isinstance(m, str) else ChatHistoryItem(**m)
            for m in session["messages"]
        ]
        result = ChatSession(
            id=str(session["id"]),
            user_id=str(session["user_id"]),
            filename=session["filename"],
            created_at=session["created_at"],
            messages=messages
        )
        logger.info(f"Fetched chat session {session_id} for user: {user.email}")
        return result
    except json.JSONDecodeError as e:
        logger.error(f"JSON parsing error in get_chat_session: {str(e)}")
        raise HTTPException(status_code=500, detail="Error parsing chat session data")
    except Exception as e:
        logger.error(f"Get chat session error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query", response_model=QueryResponse)
async def query_documents(request: QueryRequest, user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        session = await db.fetchrow(
            "SELECT id FROM chat_sessions WHERE id = $1 AND user_id = $2",
            request.session_id, user.id
        )
        if not session:
            raise HTTPException(status_code=404, detail="Chat session not found")
        
        response_message = coordinator.process_query(request.query, request.session_id)
        
        if response_message.type == "FINAL_RESPONSE":
            message_id = str(uuid.uuid4())
            await db.execute(
                """
                INSERT INTO messages (id, session_id, query, answer, sources, timestamp)
                VALUES ($1, $2, $3, $4, $5, $6)
                """,
                message_id, request.session_id, request.query, response_message.payload["answer"],
                response_message.payload["sources"], datetime.utcnow()
            )
            
            logger.info(f"Query processed for user: {user.email}, session_id: {request.session_id}, query: {request.query}")
            return QueryResponse(
                answer=response_message.payload["answer"],
                sources=response_message.payload["sources"],
                query=request.query,
                session_id=request.session_id
            )
        else:
            raise HTTPException(status_code=500, detail=response_message.payload["error"])
    except Exception as e:
        logger.error(f"Query error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))

@app.delete("/chat-history")
async def clear_chat_history(user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        await db.execute("DELETE FROM messages WHERE session_id IN (SELECT id FROM chat_sessions WHERE user_id = $1)", user.id)
        await db.execute("DELETE FROM chat_sessions WHERE user_id = $1", user.id)
        logger.info(f"Chat history cleared for user: {user.email}")
        return {"message": "Chat history cleared"}
    except Exception as e:
        logger.error(f"Clear chat history error: {str(e)}")
        raise HTTPException(status_code=400, detail=str(e))

@app.get("/stats")
async def get_stats(user: UserResponse = Depends(get_current_user), db: asyncpg.Connection = Depends(get_db)):
    try:
        sessions = await db.fetch("SELECT filename FROM chat_sessions WHERE user_id = $1", user.id)
        messages = await db.fetch("SELECT id FROM messages WHERE session_id IN (SELECT id FROM chat_sessions WHERE user_id = $1)", user.id)
        stats = {
            "total_documents": len(set([s["filename"] for s in sessions if s["filename"] != "No file uploaded"])),
            "total_chunks": len(chunks),
            "chat_history_count": len(messages)
        }
        logger.info(f"Stats retrieved for user: {user.email}, stats: {stats}")
        return stats
    except Exception as e:
        logger.error(f"Stats error: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)